"""Multi-format parser selection and canonical-document tests."""

from __future__ import annotations

from types import SimpleNamespace

import pytest

import app.rag.parsing.parsers as parsers
import app.rag.parsing.registry as registry
from app.core.config import settings
from app.rag.cleaning import EmptyContentError
from app.rag.documents import ParsedDocument, ParsedPage
from app.rag.parsing.base import TIER_FIRST_CLASS, TIER_LIGHTWEIGHT


def _select(monkeypatch, *, provider: str, docling: bool, llama: bool = False):
    monkeypatch.setattr(settings, "PARSER_PROVIDER", provider)
    monkeypatch.setattr(registry, "_docling_available", lambda: docling)
    monkeypatch.setattr(registry, "_has_llama_cloud", lambda: llama)


def _ids(ext: str) -> list[str]:
    return [parser.id for parser in registry._candidates(ext)]


def test_docling_is_default_primary_with_one_format_fallback(monkeypatch):
    _select(monkeypatch, provider="docling", docling=True)
    assert _ids(".pdf") == ["docling", "pymupdf"]
    assert _ids(".docx") == ["docling", "python_docx"]
    assert _ids(".pptx") == ["docling", "python_pptx"]
    assert _ids(".html") == ["docling", "beautifulsoup"]
    assert _ids(".xlsx") == ["openpyxl"]
    assert _ids(".txt") == ["text"]


def test_llamaparse_is_explicit_and_falls_back_locally(monkeypatch):
    _select(monkeypatch, provider="llamaparse", docling=True, llama=True)
    assert _ids(".pdf") == ["llamaparse", "docling", "pymupdf"]
    assert _ids(".docx") == ["llamaparse", "docling", "python_docx"]
    assert _ids(".html") == ["docling", "beautifulsoup"]


def test_lightweight_profile_never_loads_cloud_or_docling(monkeypatch):
    _select(monkeypatch, provider="lightweight", docling=True, llama=True)
    assert _ids(".pdf") == ["pymupdf"]
    assert _ids(".json") == ["text"]
    assert _ids(".png") == []


class _FakeParser:
    def __init__(self, parser_id, tier, result=None, error=None):
        self.id = parser_id
        self.tier = tier
        self.result = result
        self.error = error

    def supports(self, _ext):
        return True

    def parse(self, _path):
        if self.error:
            raise self.error
        return self.result


def _parsed(text: str, parser_id: str = "fake") -> ParsedDocument:
    return ParsedDocument(
        pages=[ParsedPage(text=text, number=1)],
        parser_id=parser_id,
        content_kind="markdown",
    )


def test_registry_returns_canonical_document_and_records_fallback():
    candidates = [
        _FakeParser("bad", TIER_FIRST_CLASS, error=RuntimeError("broken")),
        _FakeParser("good", TIER_LIGHTWEIGHT, result=_parsed("# Title\nbody", "good")),
    ]
    result = registry._run_candidates("unused.pdf", candidates)
    assert result is not None
    assert result.text == "# Title\nbody"
    assert result.parser_id == "good"
    assert result.parser_profile["fallback_used"] is True
    assert result.parser_profile["page_count"] == 1
    assert any("bad" in warning for warning in result.parser_profile["warnings"])


def test_registry_rejects_empty_results():
    candidates = [
        _FakeParser("blank", TIER_LIGHTWEIGHT, result=_parsed("   ", "blank"))
    ]
    assert registry._run_candidates("unused.txt", candidates) is None


def test_parse_document_raises_friendly_error_when_every_parser_fails(monkeypatch):
    monkeypatch.setattr(registry, "_candidates", lambda _ext: [])
    with pytest.raises(EmptyContentError, match="文档解析失败"):
        registry.parse_document("missing.png")


def test_docx_parser_extracts_paragraphs(tmp_path):
    import docx

    path = tmp_path / "sample.docx"
    document = docx.Document()
    document.add_heading("Caching", 1)
    document.add_paragraph("Redis avoids repeated database reads.")
    document.save(path)

    result = parsers.DocxParser().parse(str(path))
    assert "Caching" in result.pages[0].text
    assert "Redis" in result.pages[0].text


def test_pptx_parser_preserves_slide_numbers(tmp_path):
    from pptx import Presentation

    path = tmp_path / "slides.pptx"
    deck = Presentation()
    for title in ("Architecture", "Queues"):
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = title
    deck.save(path)

    result = parsers.PptxParser().parse(str(path))
    assert [page.number for page in result.pages] == [1, 2]
    assert [page.text for page in result.pages] == ["Architecture", "Queues"]


def test_xlsx_parser_makes_rows_self_describing(tmp_path):
    from openpyxl import Workbook

    path = tmp_path / "skills.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["skill", "score"])
    sheet.append(["Python", 92])
    workbook.save(path)

    result = parsers.XlsxParser().parse(str(path))
    assert result.pages[0].text == "sheet: Sheet | skill: Python | score: 92"


def test_html_parser_removes_navigation_and_keeps_structure(tmp_path):
    path = tmp_path / "page.html"
    path.write_text(
        "<nav>menu</nav><h1>Redis</h1><p>Cache data.</p><script>bad()</script>",
        encoding="utf-8",
    )
    result = parsers.HtmlParser().parse(str(path))
    assert result.content_kind == "markdown"
    assert "Redis" in result.pages[0].text
    assert "Cache data" in result.pages[0].text
    assert "menu" not in result.pages[0].text
    assert "bad()" not in result.pages[0].text


def test_text_parser_detects_non_utf8(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_bytes("缓存雪崩".encode("gb18030"))
    result = parsers.TextParser().parse(str(path))
    assert result.pages[0].text == "缓存雪崩"


def test_reader_page_labels_become_page_numbers():
    docs = [
        SimpleNamespace(text="one", metadata={"page_label": "1"}),
        SimpleNamespace(text="two", metadata={"page_label": "2"}),
    ]
    pages = parsers._pages_from_documents(docs)
    assert [(page.text, page.number) for page in pages] == [("one", 1), ("two", 2)]


def test_ocr_requires_configuration_and_runtime(monkeypatch):
    monkeypatch.setattr(settings, "RAG_OCR_ENABLED", True)
    monkeypatch.setattr(parsers, "_ocr_available", lambda: True)
    assert parsers._ocr_enabled() is True
    monkeypatch.setattr(settings, "RAG_OCR_ENABLED", False)
    assert parsers._ocr_enabled() is False


def test_docling_partial_conversion_fails_into_registry_fallback(monkeypatch):
    from docling.datamodel.base_models import ConversionStatus

    result = SimpleNamespace(
        status=ConversionStatus.PARTIAL_SUCCESS,
        errors=[SimpleNamespace(error_message="page failed")],
    )
    monkeypatch.setattr(
        parsers,
        "_get_docling_converter",
        lambda: SimpleNamespace(convert=lambda _path: result),
    )
    with pytest.raises(RuntimeError, match="partial_success"):
        parsers.DoclingParser().parse("fixture.pdf")


def test_images_and_legacy_office_are_only_claimed_by_capable_parsers():
    assert parsers.DoclingParser().supports(".png")
    assert parsers.LlamaParseParser().supports(".png")
    assert parsers.LlamaParseParser().supports(".doc")
    assert not parsers.DoclingParser().supports(".doc")
