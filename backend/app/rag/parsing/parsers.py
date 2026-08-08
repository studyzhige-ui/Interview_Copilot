"""Concrete parsers (Phase E1) — thin wrappers over the parser libraries the
project already uses, behind the :class:`DocumentParser` contract.

E1 keeps behaviour faithful to the previous inline logic: LlamaParse (first
class, when a LlamaCloud key is set) for PDF/PPTX/DOCX, PyMuPDF for PDF text,
and LlamaIndex's default ``SimpleDirectoryReader`` for everything else. Docling
(E2) and the per-format lightweight matrix (E3) register alongside these later.
"""

from __future__ import annotations

import csv
import os
from threading import Lock

from app.core.config import settings
from app.rag.documents import ParsedDocument, ParsedPage

from .base import (
    IMAGE_EXTS,
    LEGACY_OFFICE_EXTS,
    TIER_FIRST_CLASS,
    TIER_LIGHTWEIGHT,
)


def _pages_from_documents(docs: list) -> list[ParsedPage]:
    pages: list[ParsedPage] = []
    for i, d in enumerate(docs):
        text = getattr(d, "text", None) or ""
        if not text:
            continue
        meta = getattr(d, "metadata", None) or {}
        try:
            page = int(meta.get("page_label"))
        except (TypeError, ValueError):
            page = i + 1 if len(docs) > 1 else None
        pages.append(ParsedPage(text=text, number=page))
    return pages


def _read_text(file_path: str) -> str:
    """Read a text file with encoding detection (charset-normalizer), falling
    back to UTF-8 with replacement — the ingest text is always normalized to a
    str (plan §4.1.3: TXT/HTML/CSV encoding detection)."""
    from charset_normalizer import from_bytes

    with open(file_path, "rb") as fh:
        raw = fh.read()
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            pass
    best = from_bytes(raw).best()
    if best is not None:
        return str(best)
    return raw.decode("utf-8", errors="replace")


class LlamaParseParser:
    """First-class cloud parser → Markdown. Available only when a LlamaCloud key
    is configured (the registry gates on that). Claims images (cloud OCR) and
    legacy Office (.doc/.ppt/.xls, parsed directly) too — so when LlamaParse is
    the primary parser those go to the cloud instead of local RapidOCR / the
    LibreOffice conversion path (plan §4.1.3 matrix)."""

    id = "llamaparse"
    tier = TIER_FIRST_CLASS
    _EXTS = {".pdf", ".pptx", ".docx"} | IMAGE_EXTS | LEGACY_OFFICE_EXTS

    def supports(self, ext: str) -> bool:
        return ext in self._EXTS

    def parse(self, file_path: str) -> ParsedDocument:
        import nest_asyncio
        from llama_index.core import SimpleDirectoryReader
        from llama_parse import LlamaParse

        nest_asyncio.apply()
        parser = LlamaParse(
            result_type="markdown",
            language="ch_sim",
            api_key=settings.LLAMA_CLOUD_API_KEY,
            num_workers=2,
        )
        ext = os.path.splitext(file_path)[1].lower()
        docs = SimpleDirectoryReader(
            input_files=[file_path],
            file_extractor={ext: parser},
        ).load_data()
        return ParsedDocument(
            pages=_pages_from_documents(docs),
            parser_id=self.id,
            content_kind="markdown",
        )


_docling_converter = None
_docling_converter_key: tuple[str, bool] | None = None
# Docling reuses one pipeline per converter and its convert() is not thread-safe.
# The pipeline worker is solo today, but the lock keeps this module correct when
# called from tests, an API process, or a future threaded deployment.
_docling_lock = Lock()

_ocr_available_cache: bool | None = None


def _ocr_available() -> bool:
    """Whether the OCR engine (rapidocr-onnxruntime) is importable (cached).

    Docling OCR is gated on this so a deploy WITHOUT the engine still parses
    text PDFs (do_ocr=False) instead of failing on a missing engine — the same
    graceful-degradation contract the registry gives Docling itself."""
    global _ocr_available_cache
    if _ocr_available_cache is None:
        from importlib.util import find_spec

        _ocr_available_cache = find_spec("rapidocr_onnxruntime") is not None
    return _ocr_available_cache


def _ocr_enabled() -> bool:
    """OCR runs only when configured on (RAG_OCR_ENABLED) AND the engine is
    installed. On-demand within Docling: only scanned PDF pages (no text layer)
    and image documents actually OCR (plan §4.1.4: OCR 按需触发)."""
    return bool(settings.RAG_OCR_ENABLED) and _ocr_available()


def _get_docling_converter():
    """Build/return the cached Docling converter (model weights load on first
    convert). The caller holds ``_docling_lock`` — the converter is shared and
    docling's convert() isn't safe to run concurrently on it.

    OCR (RapidOCR) is wired for PDF + image inputs only when ``_ocr_enabled()``;
    otherwise the converter is built with ``do_ocr=False`` so text PDFs still
    parse without the OCR engine present. Construction is lazy (RapidOcrOptions
    is a plain options object — the engine imports at convert time, not here)."""
    global _docling_converter, _docling_converter_key
    key = (settings.RAG_DEVICE, _ocr_enabled())
    if _docling_converter is None or _docling_converter_key != key:
        # Set all model-cache environment variables before importing Docling;
        # otherwise its first import can lock in ~/.cache defaults.
        from app.core.hf_runtime import DOCLING_MODELS_DIR, prepare_hf_runtime

        prepare_hf_runtime()
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.accelerator_options import AcceleratorOptions
        from docling.datamodel.pipeline_options import (
            PdfPipelineOptions,
            RapidOcrOptions,
        )
        from docling.document_converter import (
            DocumentConverter,
            ImageFormatOption,
            PdfFormatOption,
        )

        pdf_opts = PdfPipelineOptions(
            do_ocr=_ocr_enabled(),
            ocr_options=RapidOcrOptions(),  # onnxruntime-based, deployment-light
            accelerator_options=AcceleratorOptions(device=settings.RAG_DEVICE),
        )
        # Pre-downloaded Docling artifacts share data/cache/models with the
        # embedding, reranker and speech models.
        pdf_opts.artifacts_path = DOCLING_MODELS_DIR
        _docling_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_opts),
                InputFormat.IMAGE: ImageFormatOption(pipeline_options=pdf_opts),
            }
        )
        _docling_converter_key = key
    return _docling_converter


class DoclingParser:
    """First-class LOCAL parser → Markdown (peer to LlamaParse). Available only
    when the ``docling`` package is installed (the registry gates on that and
    degrades to LlamaParse / lightweight when it isn't).

    Handles scanned PDFs and image documents via on-demand RapidOCR when
    ``RAG_OCR_ENABLED`` and the engine is installed (plan §4.1.3); without the
    engine, text PDFs still parse (do_ocr=False) and an image yields empty text
    → the registry's friendly EmptyContentError (images have no lightweight
    fallback, per the §4.1.3 matrix).

    xlsx is intentionally NOT claimed: Docling would emit a Markdown table, but
    the chunk stage routes ``.xlsx`` to the table splitter by filename (before
    the markdown branch), so xlsx stays on the lightweight path. pdf/docx/pptx/
    html and OCR'd images route to MarkdownNodeParser."""

    id = "docling"
    tier = TIER_FIRST_CLASS
    _EXTS = {".pdf", ".docx", ".pptx", ".html", ".htm"} | IMAGE_EXTS

    def supports(self, ext: str) -> bool:
        return ext in self._EXTS

    def parse(self, file_path: str) -> ParsedDocument:
        ext = os.path.splitext(file_path)[1].lower()
        with _docling_lock:  # serialize: Docling convert() isn't thread-safe
            converter = _get_docling_converter()
            result = converter.convert(file_path)
            from docling.datamodel.base_models import ConversionStatus

            if result.status is not ConversionStatus.SUCCESS:
                messages = "; ".join(
                    str(getattr(error, "error_message", None) or error)
                    for error in result.errors[:3]
                )
                raise RuntimeError(
                    f"Docling conversion was {result.status.value}"
                    + (f": {messages}" if messages else "")
                )
            document = result.document
            page_numbers = sorted(int(number) for number in document.pages)
            if page_numbers:
                from docling_core.types.doc import ContentLayer

                pages = [
                    ParsedPage(
                        text=document.export_to_markdown(
                            page_no=page_number,
                            included_content_layers={ContentLayer.BODY},
                        ),
                        number=page_number,
                    )
                    for page_number in page_numbers
                ]
            else:
                pages = [ParsedPage(text=document.export_to_markdown())]
        # ocr_used is best-effort: an image document IS OCR (its only text
        # source), so it's True when OCR is active; for text formats we don't
        # claim OCR even though Docling may OCR scanned PDF pages — there's no
        # precise per-document signal, so we never over-report (plan §4.1.4 r.10).
        ocr_used = ext in IMAGE_EXTS and _ocr_enabled()
        # page_map is empty: Docling's exported Markdown exposes no per-page char
        # spans, so page_count reads 0 here (unlike PyMuPDF, which maps per page).
        # Accurate Docling page provenance is deferred to the page_start/end round.
        return ParsedDocument(
            pages=pages,
            parser_id=self.id,
            content_kind="markdown",
            ocr_used=ocr_used,
        )


class PyMuPDFParser:
    """Lightweight PDF text extraction (no cloud, no OCR)."""

    id = "pymupdf"
    tier = TIER_LIGHTWEIGHT

    def supports(self, ext: str) -> bool:
        return ext == ".pdf"

    def parse(self, file_path: str) -> ParsedDocument:
        from llama_index.core import SimpleDirectoryReader
        from llama_index.readers.file import PyMuPDFReader

        docs = SimpleDirectoryReader(
            input_files=[file_path],
            file_extractor={".pdf": PyMuPDFReader()},
        ).load_data()
        return ParsedDocument(
            pages=_pages_from_documents(docs),
            parser_id=self.id,
            content_kind="text",
        )


# ── Per-format lightweight parsers ───────────────────────────────────────────


class DocxParser:
    """Lightweight DOCX → paragraph text (python-docx). Body paragraphs only —
    no styles, tables, headers, or footers (the first-class parsers cover those)."""

    id = "python_docx"
    tier = TIER_LIGHTWEIGHT

    def supports(self, ext: str) -> bool:
        return ext == ".docx"

    def parse(self, file_path: str) -> ParsedDocument:
        import docx

        document = docx.Document(file_path)
        text = "\n\n".join(p.text for p in document.paragraphs if p.text.strip())
        return ParsedDocument(
            pages=[ParsedPage(text=text)], parser_id=self.id, content_kind="text"
        )


class PptxParser:
    """Lightweight PPTX → per-slide text (python-pptx)."""

    id = "python_pptx"
    tier = TIER_LIGHTWEIGHT

    def supports(self, ext: str) -> bool:
        return ext == ".pptx"

    def parse(self, file_path: str) -> ParsedDocument:
        from pptx import Presentation

        slides: list[ParsedPage] = []
        for index, slide in enumerate(Presentation(file_path).slides, start=1):
            parts = [
                shape.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            if parts:
                slides.append(ParsedPage(text="\n".join(parts), number=index))
        return ParsedDocument(pages=slides, parser_id=self.id, content_kind="text")


class XlsxParser:
    """Lightweight XLSX → per-sheet self-describing rows (openpyxl). Each cell is
    emitted as ``header: value`` so a row is understandable on its own — this
    keeps multi-sheet workbooks correct even though the chunk stage currently
    routes .xlsx to the table splitter (the proper table-aware routing is E4)."""

    id = "openpyxl"
    tier = TIER_LIGHTWEIGHT

    def supports(self, ext: str) -> bool:
        return ext == ".xlsx"

    def parse(self, file_path: str) -> ParsedDocument:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        try:
            lines: list[str] = []
            for sheet in workbook.worksheets:
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    continue
                header = [str(c) if c is not None else "" for c in rows[0]]
                for row in rows[1:]:
                    cells = [
                        f"{header[i]}: {value}"
                        for i, value in enumerate(row)
                        if value is not None and i < len(header)
                    ]
                    if cells:
                        lines.append(f"sheet: {sheet.title} | " + " | ".join(cells))
        finally:
            workbook.close()
        # No sheet marker: every row already self-describes (header: value), so
        # rows from different sheets stay correct without a "[Sheet]" line that
        # the table splitter would mis-promote to a repeated header.
        return ParsedDocument(
            pages=[ParsedPage(text="\n".join(lines))],
            parser_id=self.id,
            content_kind="table",
        )


class HtmlParser:
    """Lightweight HTML → Markdown (BeautifulSoup drops script/style/nav noise,
    then markdownify keeps heading/list/table structure). Emits Markdown
    Markdown output lets the chunk stage use MarkdownNodeParser — feeding
    HTMLNodeParser tag-stripped text yields zero nodes (silent content loss)."""

    id = "beautifulsoup"
    tier = TIER_LIGHTWEIGHT

    def supports(self, ext: str) -> bool:
        return ext in (".html", ".htm")

    def parse(self, file_path: str) -> ParsedDocument:
        from bs4 import BeautifulSoup
        from markdownify import markdownify

        soup = BeautifulSoup(_read_text(file_path), "html.parser")
        # Drop only navigation/script/style noise (plan §4.1.3) — not <header>/
        # <footer>, which semantic pages often use for real body content.
        for tag in soup(["script", "style", "nav"]):
            tag.decompose()
        markdown = markdownify(str(soup)).strip()
        return ParsedDocument(
            pages=[ParsedPage(text=markdown)],
            parser_id=self.id,
            content_kind="markdown",
        )


class TextParser:
    """Normalize text, delimited tables, JSON, Markdown, and source code."""

    id = "text"
    tier = TIER_LIGHTWEIGHT
    _EXTS = {
        ".txt",
        ".csv",
        ".tsv",
        ".md",
        ".markdown",
        ".json",
        ".py",
        ".java",
        ".cpp",
        ".c",
    }

    def supports(self, ext: str) -> bool:
        return ext in self._EXTS

    def parse(self, file_path: str) -> ParsedDocument:
        ext = os.path.splitext(file_path)[1].lower()
        text = _read_text(file_path)
        if ext in {".csv", ".tsv"}:
            rows = list(
                csv.reader(
                    text.splitlines(),
                    delimiter="," if ext == ".csv" else "\t",
                )
            )
            header = rows[0] if rows else []
            records = [
                " | ".join(
                    f"{header[index]}: {value}"
                    for index, value in enumerate(row)
                    if index < len(header) and value
                )
                for row in rows[1:]
            ]
            text = "\n".join(record for record in records if record)
            content_kind = "table"
        elif ext in {".md", ".markdown"}:
            content_kind = "markdown"
        elif ext == ".json":
            content_kind = "json"
        elif ext in {".py", ".java", ".cpp", ".c"}:
            content_kind = "code"
        else:
            content_kind = "text"
        return ParsedDocument(
            pages=[ParsedPage(text=text)],
            parser_id=self.id,
            content_kind=content_kind,
        )
