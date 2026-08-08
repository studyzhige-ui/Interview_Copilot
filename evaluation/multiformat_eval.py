"""End-to-end semantic equivalence check for supported document formats.

The same small knowledge card is rendered into each format, indexed through the
production ingestion path under an isolated evaluation user, and queried through
the production hybrid retriever.  This complements the broad public-corpus
benchmark: the latter measures topic-scale retrieval quality, while this suite
detects format-specific extraction or structure loss.
"""

from __future__ import annotations

import argparse
import asyncio
import csv
import json
import sys
import tempfile
import time
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Callable

PROJECT_ROOT = Path(__file__).resolve().parents[1]
BACKEND_ROOT = PROJECT_ROOT / "backend"
if str(BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(BACKEND_ROOT))

EVAL_USER = "eval_multiformat"
TITLE = "Format Equivalence Knowledge Card"
FACTS = (
    ("Dispatch queue", "Atlas Scheduler sends urgent jobs through Quasar Queue."),
    ("Retry window", "A failed urgent job waits exactly 37 seconds before retry."),
    ("Trace header", "Every dispatched job carries the X-Orbit-Trace header."),
)
QUESTIONS = (
    ("Which queue carries urgent Atlas Scheduler jobs?", "Quasar Queue"),
    ("How long does a failed urgent Atlas job wait before retry?", "37 seconds"),
    ("Which trace header is attached to every dispatched job?", "X-Orbit-Trace"),
)


@dataclass(frozen=True)
class FormatCase:
    extension: str
    render: Callable[[Path], None]
    requires_ocr: bool = False


def _prose() -> str:
    return TITLE + "\n\n" + "\n".join(text for _label, text in FACTS)


def _write_text(path: Path) -> None:
    path.write_text(_prose(), encoding="utf-8")


def _write_markdown(path: Path) -> None:
    body = "\n\n".join(f"## {label}\n\n{text}" for label, text in FACTS)
    path.write_text(f"# {TITLE}\n\n{body}\n", encoding="utf-8")


def _write_html(path: Path) -> None:
    sections = "".join(
        f"<section><h2>{label}</h2><p>{text}</p></section>" for label, text in FACTS
    )
    path.write_text(
        f"<!doctype html><html><body><main><h1>{TITLE}</h1>{sections}</main>"
        "<nav>irrelevant navigation</nav><script>ignore()</script></body></html>",
        encoding="utf-8",
    )


def _write_json(path: Path) -> None:
    payload = {"title": TITLE, "facts": {label: text for label, text in FACTS}}
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _write_csv(path: Path) -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(("topic", "statement"))
        writer.writerows(FACTS)


def _write_python(path: Path) -> None:
    facts = json.dumps(dict(FACTS), ensure_ascii=False, indent=4)
    path.write_text(
        f'"""{TITLE}."""\n\nKNOWLEDGE = {facts}\n',
        encoding="utf-8",
    )


def _write_docx(path: Path) -> None:
    import docx

    document = docx.Document()
    document.add_heading(TITLE, 0)
    for label, text in FACTS:
        document.add_heading(label, level=1)
        document.add_paragraph(text)
    document.save(path)


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    deck = Presentation()
    for label, text in FACTS:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = label
        slide.placeholders[1].text = text
    deck.save(path)


def _write_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Knowledge"
    sheet.append(("topic", "statement"))
    for row in FACTS:
        sheet.append(row)
    workbook.save(path)


def _write_pdf(path: Path) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page()
    y = 72
    for line in (TITLE, "", *(text for _label, text in FACTS)):
        if line:
            page.insert_text((72, y), line, fontsize=11)
        y += 28
    document.save(path)
    document.close()


def _write_png(path: Path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (2200, 1250), "white")
    draw = ImageDraw.Draw(image)
    font_path = Path("C:/Windows/Fonts/arial.ttf")
    bold_path = Path("C:/Windows/Fonts/arialbd.ttf")
    body_font = ImageFont.truetype(str(font_path), 42) if font_path.is_file() else None
    heading_font = (
        ImageFont.truetype(str(bold_path), 44) if bold_path.is_file() else body_font
    )
    y = 55
    draw.text((55, y), TITLE, fill="black", font=heading_font)
    y += 145
    for label, text in FACTS:
        draw.text((55, y), label, fill="black", font=heading_font)
        draw.text((55, y + 70), text, fill="black", font=body_font)
        y += 285
    image.save(path)


FORMAT_CASES = (
    FormatCase(".txt", _write_text),
    FormatCase(".md", _write_markdown),
    FormatCase(".html", _write_html),
    FormatCase(".json", _write_json),
    FormatCase(".csv", _write_csv),
    FormatCase(".py", _write_python),
    FormatCase(".docx", _write_docx),
    FormatCase(".pptx", _write_pptx),
    FormatCase(".xlsx", _write_xlsx),
    FormatCase(".pdf", _write_pdf),
    FormatCase(".png", _write_png, requires_ocr=True),
)


def render_fixtures(directory: Path) -> list[tuple[FormatCase, Path]]:
    rendered: list[tuple[FormatCase, Path]] = []
    for case in FORMAT_CASES:
        path = directory / f"format-equivalence{case.extension}"
        case.render(path)
        rendered.append((case, path))
    return rendered


def _delete_eval_user(username: str) -> None:
    from app.db.database import SessionLocal
    from app.models.user import User
    from evaluation.prepare_corpus import _reset_user_corpus

    with SessionLocal() as db:
        user = db.query(User).filter(User.username == username).one_or_none()
        user_pk = int(user.id) if user is not None else None
    if user_pk is None:
        return
    _reset_user_corpus(user_pk)
    with SessionLocal() as db:
        user = db.query(User).filter(User.id == user_pk).one_or_none()
        if user is not None:
            db.delete(user)
            db.commit()


def _parser_metadata(document_id: str) -> dict:
    from app.db.database import SessionLocal
    from app.models.document_chunk import DocumentChunk

    with SessionLocal() as db:
        row = (
            db.query(DocumentChunk)
            .filter(DocumentChunk.document_id == document_id)
            .order_by(DocumentChunk.chunk_index)
            .first()
        )
        return json.loads(row.metadata_json or "{}") if row is not None else {}


async def _evaluate_case(case: FormatCase, path: Path, user_pk: int) -> dict:
    from app.rag.contracts import SearchIntent
    from app.rag.retriever import query_knowledge_base
    from evaluation.prepare_corpus import _document_id, _index_file, _reset_user_corpus

    _reset_user_corpus(user_pk)
    started = time.perf_counter()
    _name, chunk_count = await _index_file(path, user_pk)
    ingest_ms = (time.perf_counter() - started) * 1000
    document_id = _document_id(path)
    metadata = _parser_metadata(document_id)
    checks: list[dict] = []
    for query, expected in QUESTIONS:
        query_started = time.perf_counter()
        result = await query_knowledge_base(
            intents=[SearchIntent.from_query(query)],
            user_id=EVAL_USER,
        )
        latency_ms = (time.perf_counter() - query_started) * 1000
        context = "\n".join(str(chunk.get("text") or "") for chunk in result.chunks)
        checks.append(
            {
                "query": query,
                "expected": expected,
                "hit": result.retrieval_hit
                and expected.casefold() in context.casefold(),
                "latency_ms": round(latency_ms, 1),
                "top_score": round(float(result.chunks[0]["score"]), 4)
                if result.chunks
                else None,
                "empty_reason": result.state.empty_reason,
            }
        )
    return {
        "format": case.extension.lstrip("."),
        "parser_id": metadata.get("parser_id"),
        "fallback_used": (metadata.get("parser_profile") or {}).get("fallback_used"),
        "ocr_used": metadata.get("ocr_used", False),
        "chunk_count": chunk_count,
        "ingest_ms": round(ingest_ms, 1),
        "hits": sum(check["hit"] for check in checks),
        "queries": len(checks),
        "checks": checks,
    }


def _write_report(results: list[dict], output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    passed = sum(result["hits"] for result in results)
    total = sum(result["queries"] for result in results)
    payload = {
        "generated_at": datetime.now(UTC).isoformat(),
        "formats": len(results),
        "semantic_hits": passed,
        "semantic_queries": total,
        "semantic_hit_rate": round(passed / total, 4) if total else None,
        "results": results,
    }
    (output_dir / "report.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    lines = [
        "# Multi-format RAG equivalence",
        "",
        f"Formats: {len(results)}; semantic hits: {passed}/{total}.",
        "",
        "| Format | Parser | Fallback | OCR | Chunks | Hits | Ingest ms |",
        "|---|---|---:|---:|---:|---:|---:|",
    ]
    lines.extend(
        f"| {r['format']} | {r['parser_id']} | {r['fallback_used']} | "
        f"{r['ocr_used']} | {r['chunk_count']} | {r['hits']}/{r['queries']} | "
        f"{r['ingest_ms']} |"
        for r in results
    )
    (output_dir / "report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


async def _run(output_dir: Path) -> list[dict]:
    from app.core.config import settings
    from app.core.runtime_files import runtime_temp_dir
    from evaluation.prepare_corpus import _ensure_user
    from evaluation.runners import prepare_runtime

    prepare_runtime()
    user_pk = _ensure_user(EVAL_USER)
    previous_provider = settings.PARSER_PROVIDER
    settings.PARSER_PROVIDER = "docling"
    try:
        with tempfile.TemporaryDirectory(dir=runtime_temp_dir()) as temp:
            fixtures = render_fixtures(Path(temp))
            results = []
            for index, (case, path) in enumerate(fixtures, 1):
                print(f"[{index}/{len(fixtures)}] {case.extension}", flush=True)
                results.append(await _evaluate_case(case, path, user_pk))
    finally:
        settings.PARSER_PROVIDER = previous_provider
        _delete_eval_user(EVAL_USER)
    _write_report(results, output_dir)
    return results


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=PROJECT_ROOT / "data" / "evaluation" / "reports" / "multiformat-final",
    )
    args = parser.parse_args()
    results = asyncio.run(_run(args.output.resolve()))
    passed = sum(result["hits"] for result in results)
    total = sum(result["queries"] for result in results)
    print(f"semantic_hits={passed}/{total}")
    if passed != total:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
